import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_eureka_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Color Palette
    NAVY = RGBColor(11, 79, 159)      # #0B4F9F
    ROYAL_BLUE = RGBColor(16, 112, 224)# #1070E0
    RED = RGBColor(229, 57, 53)       # #E53935
    DARK_TEXT = RGBColor(30, 41, 59)   # #1E293B
    MUTED_TEXT = RGBColor(100, 116, 139) # #64748B
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)

    blank_layout = prs.slide_layouts[6]

    def add_header_footer(slide, title_text, team_badge="[Your Team Name]", slide_num=1):
        # Top-left Team Name Pill Badge
        badge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(2.2), Inches(0.45))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = ROYAL_BLUE
        badge_box.line.color.rgb = ROYAL_BLUE
        tf = badge_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = team_badge
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Slide Main Title
        title_box = slide.shapes.add_textbox(Inches(3.0), Inches(0.35), Inches(7.0), Inches(0.6))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY

        # Top Right Organization Label for Slides 2-6
        org_box = slide.shapes.add_textbox(Inches(9.8), Inches(0.35), Inches(3.0), Inches(0.6))
        tf_org = org_box.text_frame
        p_org = tf_org.paragraphs[0]
        p_org.text = "EUREKA 3.0!\nE-CELL ACEIT x AIC"
        p_org.font.size = Pt(10)
        p_org.font.bold = True
        p_org.font.color.rgb = MUTED_TEXT
        p_org.alignment = PP_ALIGN.RIGHT

        # Bottom Navy Footer Bar
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.333), Inches(0.6))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = NAVY
        footer_bar.line.color.rgb = NAVY
        tf_foot = footer_bar.text_frame
        tf_foot.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "@Eureka 3.0 Idea Submission Template"
        p_foot.font.size = Pt(11)
        p_foot.font.bold = True
        p_foot.font.color.rgb = WHITE
        p_foot.alignment = PP_ALIGN.CENTER

        # Slide Number
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(6.9), Inches(1.0), Inches(0.6))
        tf_num = num_box.text_frame
        tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(12)
        p_num.font.bold = True
        p_num.font.color.rgb = WHITE
        p_num.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------
    # SLIDE 1: Title Page
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Title Header
    t_box = slide1.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf1 = t_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "EUREKA 3.0! — ROAD TO ENTERPRISE 2026"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = NAVY

    p_sub = tf1.add_paragraph()
    p_sub.text = "TITLE PAGE (Idea Submission Template)"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RED

    # Main Card Box
    content_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf_c = content_box.text_frame
    tf_c.word_wrap = True

    items = [
        ("Problem Statement Title: ", "Enter Your Project / Problem Title Here"),
        ("Theme: ", "Smart Agriculture / Healthcare / FinTech / Green Energy / Open Innovation"),
        ("PS Category: ", "Software / Hardware"),
        ("Team Name: ", "Your Registered Team Name"),
        ("Team Leader Details: ", "Full Name | Roll No | Branch | Contact No | Email"),
        ("Team Members Details: ", "Member 2 to Member 7 (Full Names & University Roll Numbers)")
    ]

    for label, val in items:
        p = tf_c.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + label
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = NAVY

        run2 = p.add_run()
        run2.text = val
        run2.font.size = Pt(16)
        run2.font.color.rgb = DARK_TEXT

    footer_bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.333), Inches(0.6))
    footer_bar1.fill.solid()
    footer_bar1.fill.fore_color.rgb = NAVY
    footer_bar1.line.color.rgb = NAVY
    tf_f1 = footer_bar1.text_frame
    tf_f1.vertical_anchor = MSO_ANCHOR.MIDDLE
    pf1 = tf_f1.paragraphs[0]
    pf1.text = "@Eureka 3.0 Idea Submission Template"
    pf1.font.size = Pt(11)
    pf1.font.bold = True
    pf1.font.color.rgb = WHITE
    pf1.alignment = PP_ALIGN.CENTER

    num_box1 = slide1.shapes.add_textbox(Inches(12.2), Inches(6.9), Inches(1.0), Inches(0.6))
    tf_num1 = num_box1.text_frame
    tf_num1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_num1 = tf_num1.paragraphs[0]
    p_num1.text = "1"
    p_num1.font.size = Pt(12)
    p_num1.font.bold = True
    p_num1.font.color.rgb = WHITE
    p_num1.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------
    # SLIDE 2: Idea Title & Proposed Solution
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide2, "IDEA TITLE", "[Your Team Name]", 2)

    s2_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.2))
    tf2 = s2_box.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "❖ Proposed Solution (Describe your Idea/Solution/Prototype)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE

    bullets2 = [
        ("Detailed explanation of the proposed solution:", " Provide a clear, comprehensive description of your innovative idea, system workflow, or prototype mechanism."),
        ("How it addresses the problem:", " Explain how your solution directly eliminates market pain points and addresses target user challenges."),
        ("Innovation and uniqueness of the solution:", " Highlight key USPs, proprietary algorithms, novel hardware design, or competitive advantage.")
    ]

    for b_title, b_desc in bullets2:
        p = tf2.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_title
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = NAVY

        run2 = p.add_run()
        run2.text = b_desc
        run2.font.size = Pt(16)
        run2.font.color.rgb = DARK_TEXT

    # -------------------------------------------------------------
    # SLIDE 3: Technical Approach
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide3, "TECHNICAL APPROACH", "[Your Team Name]", 3)

    s3_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2))
    tf3 = s3_box.text_frame
    tf3.word_wrap = True

    bullets3 = [
        ("Technologies to be used:", " Mention programming languages, frameworks, hardware microcontrollers, sensors, AI/ML models, or cloud services (e.g. Python, React, ESP32, TensorFlow)."),
        ("Methodology and process for implementation:", " Illustrate flowcharts, system architecture diagrams, data flow process, or photos of your working prototype.")
    ]

    for i, (b_title, b_desc) in enumerate(bullets3):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_title
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = NAVY

        run2 = p.add_run()
        run2.text = b_desc
        run2.font.size = Pt(16)
        run2.font.color.rgb = DARK_TEXT

    # -------------------------------------------------------------
    # SLIDE 4: Feasibility and Viability
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide4, "FEASIBILITY AND VIABILITY", "[Your Team Name]", 4)

    s4_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2))
    tf4 = s4_box.text_frame
    tf4.word_wrap = True

    bullets4 = [
        ("Analysis of the feasibility of the idea:", " Evaluate practical execution, technical readiness, cost feasibility, and deployment timeline."),
        ("Potential challenges and risks:", " Identify hardware dependencies, data privacy, scaling bottlenecks, or adoption barriers."),
        ("Strategies for overcoming these challenges:", " Present risk mitigation plans, backup components, fail-safe mechanisms, or security protocols.")
    ]

    for i, (b_title, b_desc) in enumerate(bullets4):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_title
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = NAVY

        run2 = p.add_run()
        run2.text = b_desc
        run2.font.size = Pt(16)
        run2.font.color.rgb = DARK_TEXT

    # -------------------------------------------------------------
    # SLIDE 5: Impact and Benefits
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide5, "IMPACT AND BENEFITS", "[Your Team Name]", 5)

    s5_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2))
    tf5 = s5_box.text_frame
    tf5.word_wrap = True

    bullets5 = [
        ("Potential impact on the target audience:", " Describe the direct measurable impact on end users, industry sectors, or society."),
        ("Benefits of the solution:", " List social, economic, environmental, efficiency, or cost-reduction advantages.")
    ]

    for i, (b_title, b_desc) in enumerate(bullets5):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_title
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = NAVY

        run2 = p.add_run()
        run2.text = b_desc
        run2.font.size = Pt(16)
        run2.font.color.rgb = DARK_TEXT

    # -------------------------------------------------------------
    # SLIDE 6: Research and References
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide6, "RESEARCH AND REFERENCES", "[Your Team Name]", 6)

    s6_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2))
    tf6 = s6_box.text_frame
    tf6.word_wrap = True

    p = tf6.paragraphs[0]
    run1 = p.add_run()
    run1.text = "• Details / Links of the reference and research work:"
    run1.font.bold = True
    run1.font.size = Pt(18)
    run1.font.color.rgb = NAVY

    run2 = p.add_run()
    run2.text = " Include citations of research papers, market surveys, patent links, open-source libraries, or competitor analysis links."
    run2.font.size = Pt(16)
    run2.font.color.rgb = DARK_TEXT

    # Save presentation
    save_files = [
        "c:\\Users\\Krishna\\Documents\\Eureka\\Eureka_3.0_Pitch_Deck_Template_Final.pptx",
        "c:\\Users\\Krishna\\Documents\\Eureka\\Eureka_3.0_Pitch_Deck_Template_Updated.pptx",
        "c:\\Users\\Krishna\\Documents\\Eureka\\Eureka_3.0_Pitch_Deck_Template.pptx"
    ]
    for path in save_files:
        try:
            prs.save(path)
            print("Saved PPTX successfully at:", path)
        except Exception as e:
            pass

if __name__ == "__main__":
    create_eureka_ppt()
